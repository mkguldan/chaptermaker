"""
PowerPoint to image conversion using LibreOffice or alternative methods
"""

import logging
import subprocess
import platform
from pathlib import Path
import tempfile
import shutil
from typing import List, Dict, Any
from PIL import Image

logger = logging.getLogger(__name__)


class PresentationConverter:
    """Convert presentations to images using various methods"""
    
    def __init__(self):
        self.system = platform.system()
        self.libreoffice_path = self._find_libreoffice()
        
    def _find_libreoffice(self) -> str:
        """Find LibreOffice installation path"""
        if self.system == "Windows":
            paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
            ]
        elif self.system == "Linux":
            paths = [
                "/usr/bin/soffice",
                "/usr/bin/libreoffice",
                "/opt/libreoffice/program/soffice"
            ]
        else:  # macOS
            paths = [
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "/usr/local/bin/soffice"
            ]
            
        for path in paths:
            if Path(path).exists():
                logger.info(f"Found LibreOffice at: {path}")
                return path
                
        logger.warning("LibreOffice not found. PPT conversion may be limited.")
        return None
        
    async def convert_pptx_to_images(
        self,
        pptx_path: str,
        output_dir: str,
        output_format: str = "jpg"
    ) -> List[Dict[str, Any]]:
        """Convert PowerPoint to images"""
        
        # Try LibreOffice first
        if self.libreoffice_path:
            try:
                return await self._convert_with_libreoffice(
                    pptx_path, output_dir, output_format
                )
            except Exception as e:
                logger.error(f"LibreOffice conversion failed: {str(e)}")
                
        # Fallback to python-pptx method
        return await self._convert_with_python_pptx(
            pptx_path, output_dir, output_format
        )
        
    async def _convert_with_libreoffice(
        self,
        pptx_path: str,
        output_dir: str,
        output_format: str
    ) -> List[Dict[str, Any]]:
        """Convert using LibreOffice headless mode"""
        
        # Create temporary directory for PDF output
        temp_dir = tempfile.mkdtemp()
        
        try:
            # Convert PPT to PDF first
            cmd = [
                self.libreoffice_path,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", temp_dir,
                pptx_path
            ]
            
            logger.info(f"Running LibreOffice conversion: {' '.join(cmd)}")
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                check=True
            )
            
            # Find the generated PDF
            pdf_files = list(Path(temp_dir).glob("*.pdf"))
            if not pdf_files:
                raise ValueError("No PDF generated from PowerPoint")
                
            pdf_path = str(pdf_files[0])
            
            # Convert PDF to images
            from pdf2image import convert_from_path
            
            images = convert_from_path(
                pdf_path,
                dpi=300,
                fmt=output_format
            )
            
            # Save images
            slides = []
            # Map format for PIL (jpg -> JPEG)
            pil_format = "JPEG" if output_format.lower() == "jpg" else output_format.upper()
            
            for idx, image in enumerate(images, 1):
                filename = f"{idx:02d}.{output_format}"
                output_path = Path(output_dir) / filename
                
                # Ensure output directory exists
                output_path.parent.mkdir(parents=True, exist_ok=True)
                
                # Save image with high quality
                image.save(output_path, pil_format, quality=95)
                
                slides.append({
                    "number": idx,
                    "local_path": str(output_path),
                    "filename": filename
                })
                
            logger.info(f"Converted {len(slides)} slides using LibreOffice")
            return slides
            
        finally:
            # Clean up temp directory
            if Path(temp_dir).exists():
                shutil.rmtree(temp_dir, ignore_errors=True)
                
    async def _convert_with_python_pptx(
        self,
        pptx_path: str,
        output_dir: str,
        output_format: str
    ) -> List[Dict[str, Any]]:
        """Convert using python-pptx (basic method)"""
        
        from pptx import Presentation
        from pptx.util import Inches
        import io
        
        prs = Presentation(pptx_path)
        slides = []
        
        # Ensure output directory exists
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        
        for idx, slide in enumerate(prs.slides, 1):
            # Create a high-resolution image for the slide
            # Note: This is a basic implementation
            
            # Get slide dimensions
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Create image with white background
            img_width = 1920
            img_height = int(img_width * slide_height / slide_width)
            
            img = Image.new('RGB', (img_width, img_height), 'white')
            
            # For a more complete implementation, we would need to:
            # 1. Extract text from shapes and text boxes
            # 2. Extract images from the slide
            # 3. Recreate the layout
            # This is a placeholder that creates a white slide
            
            # Add slide number as text (placeholder)
            from PIL import ImageDraw, ImageFont
            draw = ImageDraw.Draw(img)
            
            try:
                # Try to use a nice font
                font = ImageFont.truetype("arial.ttf", 60)
            except:
                # Fallback to default font
                font = ImageFont.load_default()
                
            # Draw slide number
            text = f"Slide {idx}"
            draw.text((50, 50), text, fill='black', font=font)
            
            # Extract title if available
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                draw.text((50, 150), title_text[:50], fill='black', font=font)
            
            # Save image
            filename = f"{idx:02d}.{output_format}"
            output_path = Path(output_dir) / filename
            
            img.save(output_path, output_format.upper(), quality=95)
            
            slides.append({
                "number": idx,
                "local_path": str(output_path),
                "filename": filename
            })
            
        logger.warning(
            f"Used basic python-pptx conversion for {len(slides)} slides. "
            "For better quality, install LibreOffice."
        )
        
        return slides
